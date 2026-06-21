from app import client
import os
import concurrent.futures
from tqdm import tqdm
import uuid
import time
from typing import Any, Type, TypeVar, Union
from pydantic import BaseModel, ValidationError
from openai import OpenAI, APIError
from celery.utils.log import get_task_logger
import pandas as pd
import base64
import mimetypes
import json
import re
import time
import logging


logger = logging.getLogger(__name__)

T = TypeVar('T', bound=BaseModel)


def create_vector_store(store_name: str) -> dict:
    try:
        vector_store = client.vector_stores.create(
            name=f"{store_name}_{str(uuid.uuid4())[:8]}"
        )

        logger.info(f"Vector store type: {type(vector_store)}")
        logger.info(f"Vector store attributes: {dir(vector_store)}")
        logger.info(f"Vector store content: {vector_store}")

        if hasattr(vector_store, "id"):
            logger.info("Has 'id' attribute")
            vector_id = vector_store.id
        elif hasattr(vector_store, "get"):
            logger.info("Has 'get' method - treating as dict-like")
            vector_id = vector_store.get("id")
        else:
            logger.info("Converting to dict")
            vector_dict = (
                dict(vector_store)
                if hasattr(vector_store, "__iter__")
                else str(vector_store)
            )
            logger.info(f"Dict conversion: {vector_dict}")
            return {}

        details = {
            "id": vector_id,
            "name": getattr(vector_store, "name", "unknown"),
            "created_at": getattr(vector_store, "created_at", "unknown"),
            "file_count": getattr(
                getattr(vector_store, "file_counts", {}), "completed", 0
            ),
        }

        logger.info(f"Vector store created: {details}")
        return details

    except Exception as e:
        logger.warning(f"Error creating vector store: {e}")
        logger.warning(f"Error type: {type(e)}")
        import traceback
        logger.warning(f"Traceback: {traceback.format_exc()}")
        return {}


def upload_single_file(file_path: str, vector_store_id: str):
    try:
        with open(file_path, "rb") as f:
            file_response = client.files.create(file=f, purpose="assistants")
            print(file_response)

        client.vector_stores.files.create(
            vector_store_id=vector_store_id, file_id=file_response.id
        )

        print(f"Uploaded {os.path.basename(file_path)} to vector store {vector_store_id}")
        return {"status": "success", "file_id": file_response.id}
    except Exception as e:
        print(f"Error with {file_path}: {str(e)}")
        return {"status": "failed", "error": str(e)}


def delete_vector_store(vector_store_id: str) -> dict:
    if not client:
        return {"status": "failed", "error": "OpenAI client not initialized."}
    try:
        logger.info(f"Attempting to delete vector store with ID: {vector_store_id}")
        response = client.vector_stores.delete(vector_store_id=vector_store_id)
        logger.info(f"Successfully deleted vector store. Response: {response}")
        return {"status": "success", "id": response.id, "deleted": response.deleted}
    except Exception as e:
        logger.error(f"Error deleting vector store {vector_store_id}: {e}", exc_info=True)
        return {"status": "failed", "error": str(e)}


def delete_file(file_id: str) -> dict:
    if not client:
        return {"status": "failed", "error": "OpenAI client not initialized."}
    try:
        logger.info(f"Attempting to delete file with ID: {file_id}")
        response = client.files.delete(file_id=file_id)
        logger.info(f"Successfully deleted file. Response: {response}")
        return {"status": "success", "id": response.id, "deleted": response.deleted}
    except Exception as e:
        logger.error(f"Error deleting file {file_id}: {e}", exc_info=True)
        return {"status": "failed", "error": str(e)}


def extract_structured_info(
    query: str,
    vector_store_id: str,
    schema: Any,
    retries: int = 2,
    backoff_factor: float = 1.5,
) -> Any | None:
    """
    Extracts structured info using chat completions — provider agnostic.
    vector_store_id parameter kept for backward compatibility but ignored.
    Works with Azure OpenAI, OpenAI, and any chat completions compatible provider.
    """
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4.1-mini")
    attempt = 0
    while attempt < retries:
        try:
            logger.info("Attempt #%d to extract structured info for schema: %s", attempt + 1, schema.__name__)
            response = client.chat.completions.create(
                model=deployment,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are an expert compliance consultant. "
                            "Return ONLY valid JSON that matches the requested schema. "
                            "No markdown, no explanation, just JSON."
                        ),
                    },
                    {"role": "user", "content": query},
                ],
                response_format={"type": "json_object"},
                temperature=0.0,
            )
            raw_json = response.choices[0].message.content
            if not raw_json:
                raise ValueError("Model returned empty response.")
            raw_json = raw_json.strip()
            if raw_json.startswith("```"):
                raw_json = re.sub(r"^```(?:json)?\n?", "", raw_json)
                raw_json = re.sub(r"\n?```$", "", raw_json).strip()
            parsed = schema.model_validate_json(raw_json)
            logger.info("Successfully extracted structured data.")
            return parsed
        except (ValidationError, ValueError) as e:
            logger.warning("Attempt #%d failed with validation error: %s", attempt + 1, e)
            if attempt == retries - 1 and schema.__name__ == "ComplianceRequirements":
                logger.info("Trying regex fallback for compliance activities...")
                return _extract_compliance_with_regex_fallback(query, [])
        except Exception as e:
            logger.warning("Attempt #%d failed with general error: %s", attempt + 1, e)
        attempt += 1
        if attempt < retries:
            time.sleep(backoff_factor ** attempt)
    logger.error("All %d attempts failed. Returning None.", retries)
    return None

def _extract_compliance_with_regex_fallback(query: str, tools: list) -> Any | None:
    try:
        import re
        import json

        if tools:
            response = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {
                        "role": "system",
                        "content": "You are a helpful AI that extracts structured information from regulatory documents.",
                    },
                    {"role": "user", "content": query},
                ],
            )
        else:
            response = client.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[
                    {
                        "role": "system",
                        "content": "You are a helpful AI that extracts structured information from regulatory documents.",
                    },
                    {"role": "user", "content": query},
                ],
            )

        result = response.choices[0].message.content
        logger.info("Regex fallback - Raw AI response: %s", result[:500] if result else "Empty")

        pattern = r"```json(.*?)```"
        matches = re.findall(pattern, result, re.DOTALL)

        if matches:
            json_str = matches[0].strip()
            logger.info("Regex fallback - Extracted JSON: %s", json_str[:500])

            data = json.loads(json_str)
            transformed_data = _transform_compliance_data(data, query)

            if transformed_data and "compliance_activities" in transformed_data:
                from app.services.prompt_templates.compliance_activity import (
                    ComplianceRequirements,
                )
                return ComplianceRequirements(**transformed_data)

        logger.warning("Regex fallback also failed to extract valid JSON")
        return None

    except Exception as e:
        logger.error("Regex fallback failed: %s", str(e))
        return None


def _transform_compliance_data(raw_data: dict, query: str) -> dict:
    try:
        clause_text = _extract_clause_from_query(query)
        activities = []

        if "compliance_activities" in raw_data:
            activities = raw_data["compliance_activities"]
        elif "ComplianceRequirements" in raw_data:
            activities = raw_data["ComplianceRequirements"]
        elif isinstance(raw_data, list):
            activities = raw_data
        else:
            for key, value in raw_data.items():
                if isinstance(value, list) and value and isinstance(value[0], dict):
                    activities = value
                    break

        if not activities:
            return None

        transformed_activities = []

        for i, activity in enumerate(activities, 1):
            transformed_activity = {}

            relevant_depts = activity.get("relevant_departments", "Unknown")
            if isinstance(relevant_depts, list):
                transformed_activity["relevant_departments"] = ", ".join(relevant_depts)
            else:
                transformed_activity["relevant_departments"] = str(relevant_depts)

            activity_id = activity.get("activity_id", i)
            if isinstance(activity_id, int):
                transformed_activity["activity_id"] = str(activity_id)
            else:
                transformed_activity["activity_id"] = str(activity_id)

            clause = activity.get("clause", clause_text)
            transformed_activity["clause"] = clause

            transformed_activity["compliance_level"] = str(activity.get("compliance_level", "Design"))
            transformed_activity["department_id"] = int(activity.get("department_id", 0))
            transformed_activity["process_name"] = str(activity.get("process_name", "Unknown"))
            transformed_activity["sub_process_name"] = str(activity.get("sub_process_name", "Unknown"))
            transformed_activity["activity_description"] = str(activity.get("activity_description", ""))
            transformed_activity["responsible_party"] = str(activity.get("responsible_party", "Unknown"))
            transformed_activity["frequency"] = str(activity.get("frequency", "As needed"))
            transformed_activity["evidence_required"] = str(activity.get("evidence_required", ""))
            transformed_activity["justification"] = str(activity.get("justification", ""))

            transformed_activities.append(transformed_activity)

        return {"compliance_activities": transformed_activities}

    except Exception as e:
        logger.error("Data transformation failed: %s", str(e))
        return None


def _extract_clause_from_query(query: str) -> str:
    try:
        import re
        match = re.search(r"Regulatory Clause:\s*(.*?)(?:\n|$)", query)
        if match:
            return match.group(1).strip()
        return "Clause text not available"
    except:
        return "Clause text not available"


def extract_structured_info_2(
    query: str, schema: Any, retries: int = 2, backoff_factor: float = 1.5
) -> Any | None:
    """
    Extracts structured info using chat completions — provider agnostic.
    Works with Azure OpenAI, OpenAI, and any chat completions compatible provider.
    """
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4.1-mini")
    attempt = 0
    while attempt < retries:
        try:
            logger.info("Attempt #%d to extract structured info for schema: %s", attempt + 1, schema.__name__)
            response = client.chat.completions.create(
                model=deployment,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are an expert compliance consultant. "
                            "Return ONLY valid JSON that matches the requested schema. "
                            "No markdown, no explanation, just JSON."
                        ),
                    },
                    {"role": "user", "content": query},
                ],
                response_format={"type": "json_object"},
                temperature=0.0,
                max_tokens=16000,
            )
            raw_json = response.choices[0].message.content
            if not raw_json:
                raise ValueError("Model returned empty response.")
            raw_json = raw_json.strip()
            if raw_json.startswith("```"):
                raw_json = re.sub(r"^```(?:json)?\n?", "", raw_json)
                raw_json = re.sub(r"\n?```$", "", raw_json).strip()
            parsed = schema.model_validate_json(raw_json)
            logger.info("Successfully extracted structured data.")
            return parsed
        except (ValidationError, ValueError) as e:
            logger.warning("Attempt #%d failed with validation error: %s", attempt + 1, e)
        except Exception as e:
            logger.warning("Attempt #%d failed with general error: %s", attempt + 1, e)
        attempt += 1
        if attempt < retries:
            time.sleep(backoff_factor ** attempt)
    logger.error("All %d attempts failed. Returning default.", retries)
    if issubclass(schema, BaseModel):
        try:
            return schema.model_construct()
        except Exception as e:
            logger.error("Could not construct empty schema: %s", e)
            return None
    return {}

    try:
        if not os.path.exists(image_path):
            logger.error(f"Image file not found: {image_path}")
            raise FileNotFoundError(f"Image file not found: {image_path}")

        mime_type, _ = mimetypes.guess_type(image_path)
        if not mime_type or not mime_type.startswith("image"):
            logger.error(f"Invalid or unknown image file type for {image_path}")
            raise ValueError(f"Invalid image file type: {mime_type}")

        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
            result = f"{encoded_string}"
            logger.info(f"Successfully encoded image {image_path}, size: {len(result)} chars")
            return result
    except Exception as e:
        logger.error(f"Failed to encode image {image_path}: {e}")
        raise


def _get_config(file_name: str | None) -> dict:
    config = {
        "model": "gpt-4.1-mini",
        "preprocess": "text",
    }

    if not file_name:
        return config

    _, extension = os.path.splitext(file_name)
    extension = extension.lower()

    CONFIG_BY_EXTENSION = {
        ".pdf":  {"model": "gpt-4.1-mini", "preprocess": "pdf"},
        ".doc":  {"model": "gpt-4.1-mini", "preprocess": "docx"},
        ".docx": {"model": "gpt-4.1-mini", "preprocess": "docx"},
        ".ppt":  {"model": "gpt-4.1-mini", "preprocess": "ppt"},
        ".pptx": {"model": "gpt-4.1-mini", "preprocess": "ppt"},
        ".xls":  {"model": "gpt-4.1-mini", "preprocess": "excel"},
        ".xlsx": {"model": "gpt-4.1-mini", "preprocess": "excel"},
        ".csv":  {"model": "gpt-4.1-mini", "preprocess": "csv"},
        ".png":  {"model": "gpt-4.1-mini", "preprocess": "image"},
        ".jpg":  {"model": "gpt-4.1-mini", "preprocess": "image"},
        ".jpeg": {"model": "gpt-4.1-mini", "preprocess": "image"},
        ".gif":  {"model": "gpt-4.1-mini", "preprocess": "image"},
        ".bmp":  {"model": "gpt-4.1-mini", "preprocess": "image"},
        ".webp": {"model": "gpt-4.1-mini", "preprocess": "image"},
        ".mp3":  {"model": "gpt-4.1-mini", "preprocess": "audio"},
        ".wav":  {"model": "gpt-4.1-mini", "preprocess": "audio"},
        ".m4a":  {"model": "gpt-4.1-mini", "preprocess": "audio"},
        ".flac": {"model": "gpt-4.1-mini", "preprocess": "audio"},
        ".ogg":  {"model": "gpt-4.1-mini", "preprocess": "audio"},
        ".aac":  {"model": "gpt-4.1-mini", "preprocess": "audio"},
    }

    return CONFIG_BY_EXTENSION.get(extension, config)


def _preprocess_file(file_name: str, preprocess_type: str) -> str:
    try:
        if preprocess_type == "excel":
            df = pd.read_excel(file_name)
            return df.to_csv(index=False)

        if preprocess_type == "csv":
            df = pd.read_csv(file_name)
            return df.to_csv(index=False)

        if preprocess_type == "audio":
            if not os.path.exists(file_name):
                raise FileNotFoundError(f"Audio file not found: {file_name}")

            with open(file_name, "rb") as audio_file:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1", file=audio_file
                )
            print("Transcrib", transcript)
            if not transcript or not transcript.text.strip():
                raise ValueError(f"Audio transcription returned empty result for {file_name}")

            logger.info(f"Successfully transcribed audio file: {file_name}, text length: {len(transcript.text)}")
            return transcript.text

    except Exception as e:
        logger.error("Failed to preprocess %s file '%s': %s", preprocess_type, file_name, e)

    return ""


def extract_structured_info_3(
    query: str,
    schema: Any,
    file_name: str | None = None,
    vector_store_id: str | None = None,  # Keep parameter but don't use it
    retries: int = 3,
    backoff_factor: float = 1.5,
) -> Any | None:
    """
    Extracts structured info from a model WITHOUT vector store.
    """
    config = _get_config(file_name)
    model_name = config["model"]

    attempt = 0
    while attempt < retries:
        try:
            api_params = {
                "model": model_name,
                "text_format": schema,
            }

            # Handle different file types
            if config["preprocess"] in {"excel", "csv", "audio"} and file_name:
                extracted_text = _preprocess_file(file_name, config["preprocess"])
                if not extracted_text:
                    raise ValueError(f"Preprocessing failed for file {file_name}")
                api_params["input"] = (
                    f"{query}\n\nHere is the content from the provided file:\n{extracted_text}"
                )

            elif config["preprocess"] == "image" and file_name:
                base64_image = encode_image_to_base64(file_name)
                if not base64_image:
                    raise ValueError(f"Could not encode image to base64: {file_name}")
                api_params["input"] = [
                    {
                        "role": "user",
                        "content": [
                            {"type": "input_text", "text": query},
                            {
                                "type": "input_image",
                                "image_url": f"data:image/jpeg;base64,{base64_image}",
                            },
                        ],
                    }
                ]

            else:
                # FIX: Extract text from file and inject into prompt
                # Previously only query was sent — document content was never reaching the LLM
                if file_name:
                    try:
                        extracted_text = ""
                        ext = os.path.splitext(file_name)[1].lower()

                        if ext == ".pdf":
                            import pdfplumber
                            with pdfplumber.open(file_name) as pdf:
                                extracted_text = "\n".join(
                                    page.extract_text() or "" for page in pdf.pages
                                )
                        elif ext == ".docx":
                            from docx import Document as DocxDocument
                            doc = DocxDocument(file_name)
                            extracted_text = "\n".join(
                                p.text for p in doc.paragraphs if p.text.strip()
                            )
                        elif ext in {".ppt", ".pptx"}:
                            from pptx import Presentation
                            prs = Presentation(file_name)
                            lines = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        lines.append(shape.text.strip())
                            extracted_text = "\n".join(lines)
                        else:
                            extracted_text = _preprocess_file(file_name, "text")

                        if extracted_text and extracted_text.strip():
                            api_params["input"] = (
                                f"{query}\n\nDOCUMENT CONTENT:\n{extracted_text}"
                            )
                        else:
                            logger.warning(f"No text extracted from {file_name}, sending query only")
                            api_params["input"] = query

                    except Exception as e:
                        logger.error(f"Text extraction failed for {file_name}: {e}")
                        api_params["input"] = query
                else:
                    api_params["input"] = query

            # Make API call
            print(api_params)
            response = client.responses.parse(**api_params)
            print("response", response)
            if response and response.output_parsed is not None:
                logger.info("Successfully parsed structured data.")
                return response.output_parsed

        except (ValidationError, ValueError) as e:
            logger.warning("Validation/model error: %s", e)
            raise ValueError("Model returned null/empty response.")
        except APIError as e:
            logger.warning("API error: %s", e)
        except Exception as e:
            logger.warning("Unexpected error: %s", e)

        attempt += 1
        if attempt < retries:
            wait_time = backoff_factor**attempt
            logger.info("Retrying in %.2f seconds...", wait_time)
            time.sleep(wait_time)

    if issubclass(schema, BaseModel):
        try:
            return schema.model_construct()
        except Exception as e:
            logger.error("Could not construct empty schema instance: %s", e)
            return None

    return None
